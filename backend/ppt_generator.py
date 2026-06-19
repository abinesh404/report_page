import os
import io
from datetime import datetime, date
import psycopg2
from psycopg2.extras import RealDictCursor
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import matplotlib
matplotlib.use('Agg') # Ensure headless compatibility
import matplotlib.pyplot as plt
import copy
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
import pandas as pd
import numpy as np

def format_date_long(val):
    if not val:
        return "N/A"
    if isinstance(val, (datetime, date)):
        return val.strftime("%B %d, %Y")
    try:
        dt = datetime.strptime(str(val).split('T')[0], "%Y-%m-%d")
        return dt.strftime("%B %d, %Y")
    except Exception:
        return str(val)

def replace_text_in_shape(shape, search_text, replace_text):
    if not shape.has_text_frame:
        return False
    replaced = False
    for paragraph in shape.text_frame.paragraphs:
        if search_text in paragraph.text:
            for run in paragraph.runs:
                if search_text in run.text:
                    run.text = run.text.replace(search_text, replace_text)
                    replaced = True
                    break
            else:
                # Merge runs and replace
                cur_text = paragraph.text
                new_text = cur_text.replace(search_text, replace_text)
                if paragraph.runs:
                    first_run = paragraph.runs[0]
                    font_name = first_run.font.name
                    font_size = first_run.font.size
                    bold = first_run.font.bold
                    color = first_run.font.color.rgb if first_run.font.color else None
                    
                    for i in range(len(paragraph.runs)):
                        paragraph.runs[i].text = ""
                    first_run.text = new_text
                    first_run.font.name = font_name
                    first_run.font.size = font_size
                    first_run.font.bold = bold
                    if color:
                        first_run.font.color.rgb = color
                else:
                    paragraph.add_run().text = new_text
                replaced = True
    return replaced

def add_cell_text(slide, left, top, width, height, text, font_size_pt=10, bold=False, color_rgb=(0, 0, 0), align_left=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    run.font.name = "Inter"
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color_rgb)
    if not align_left:
        p.alignment = 1 # Center

def duplicate_slide(prs, source_slide):
    target_slide = prs.slides.add_slide(source_slide.slide_layout)
    for shp in list(target_slide.shapes):
        sp = shp._element
        sp.getparent().remove(sp)
    for shape in source_slide.shapes:
        if getattr(shape, "has_table", False):
            continue
        if shape.shape_type == 13: # PICTURE
            try:
                image_bytes = shape.image.blob
                image_stream = io.BytesIO(image_bytes)
                target_slide.shapes.add_picture(image_stream, shape.left, shape.top, shape.width, shape.height)
            except Exception:
                pass
        else:
            new_el = copy.deepcopy(shape._element)
            target_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return target_slide

def generate_presentation(report_name, sentinel_pool, db_config_complibear, report_type=None):
    # 1. Fetch data from sentinel_db (audit_plan)
    conn_sent = sentinel_pool.getconn()
    try:
        with conn_sent.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("SET search_path TO sentinel_db;")
            cur.execute("""
                SELECT id, title, risk_level, "rick score" as risk_score, process, auditor, company, sector, start_date, end_date, control_description 
                FROM audit_plan;
            """)
            controls = cur.fetchall()
            
            # Fetch dynamic company, sector, and date ranges from PostgreSQL
            cur.execute("""
                SELECT company, sector, MIN(start_date) as min_start, MAX(end_date) as max_end 
                FROM audit_plan 
                GROUP BY company, sector 
                LIMIT 1;
            """)
            meta = cur.fetchone()

            # Fetch detailed risk descriptions for detailed report
            cur.execute("""
                SELECT "Risk number", "Risk Description" 
                FROM input_rcm;
            """)
            rcm_risks = {r['Risk number']: r['Risk Description'] for r in cur.fetchall() if r['Risk number'] and r['Risk Description']}

            # Fetch insights titles & descriptions
            cur.execute("""
                SELECT title, description 
                FROM insights_current;
            """)
            insights_list = [(r['title'], r['description']) for r in cur.fetchall() if r['title'] and r['description']]
    finally:
        sentinel_pool.putconn(conn_sent)

    if not meta:
        meta = {
            "company": "CJSJ",
            "sector": "Manufacturing",
            "min_start": "2026-01-01",
            "max_end": "2026-09-30"
        }

    company = meta.get("company", "CJSJ")
    sector = meta.get("sector", "Manufacturing")
    db_start_date = format_date_long(meta.get("min_start"))
    db_end_date = format_date_long(meta.get("max_end"))

    # 2. Fetch exception counts from complibear database
    conn_comp = psycopg2.connect(**db_config_complibear)
    try:
        with conn_comp.cursor() as cur:
            cur.execute("SET search_path TO complibear;")
            tables = [
                "bank_account_changed", "cjs1_quality_rejected", "cjsa22_foc_discount",
                "cjsa23_sales_return_qty", "direct_changes_sap", "duplicate_customers",
                "finished_goods_dispatched_wo_qi", "gst_working", "mjot06_yield_loss",
                "multiple_sales_return", "password_test", "po_split", "po_terms_changed",
                "procurement_higher_contract", "reorder_level", "sales_return_180",
                "sales_return_im", "sales_return_price_mismatch", "scrap_sales",
                "sjin13_unplanned_delivery", "sjpa7_msme_penalty", "tds_insight"
            ]
            
            table_to_process = {
                "bank_account_changed": "Procure to Pay",
                "po_split": "Procure to Pay",
                "po_terms_changed": "Procure to Pay",
                "procurement_higher_contract": "Procure to Pay",
                "sjin13_unplanned_delivery": "Procure to Pay",
                "sjpa7_msme_penalty": "Procure to Pay",
                
                "cjs1_quality_rejected": "Quality Management",
                "finished_goods_dispatched_wo_qi": "Quality Management",
                
                "cjsa22_foc_discount": "Order to Cash",
                "cjsa23_sales_return_qty": "Order to Cash",
                "duplicate_customers": "Order to Cash",
                "multiple_sales_return": "Order to Cash",
                "sales_return_180": "Order to Cash",
                "sales_return_im": "Order to Cash",
                "sales_return_price_mismatch": "Order to Cash",
                
                "direct_changes_sap": "Cyber Security",
                "password_test": "Cyber Security",
                
                "gst_working": "Taxation",
                "tds_insight": "Taxation",
                
                "mjot06_yield_loss": "Inventory Management",
                "reorder_level": "Inventory Management",
                
                "scrap_sales": "Scrap Management"
            }
            
            # Fetch existing tables to prevent UndefinedTable exceptions
            cur.execute("SELECT table_name FROM information_schema.tables WHERE table_schema = 'complibear';")
            existing_tables = {r[0] for r in cur.fetchall()}

            process_counts = {p: 0 for p in set(table_to_process.values())}
            for table in tables:
                if table in existing_tables:
                    cur.execute(f"SELECT COUNT(*) FROM \"{table}\";")
                    count = cur.fetchone()[0]
                else:
                    count = 0
                process = table_to_process[table]
                process_counts[process] += count
    finally:
        conn_comp.close()

    # 3. Calculate statistics for Risk Distribution
    high_count = sum(1 for c in controls if c['risk_level'].upper() == 'HIGH')
    medium_count = sum(1 for c in controls if c['risk_level'].upper() == 'MEDIUM')
    low_count = sum(1 for c in controls if c['risk_level'].upper() == 'LOW')
    total_count = len(controls) if len(controls) > 0 else 1

    high_percent = (high_count / total_count) * 100
    medium_percent = (medium_count / total_count) * 100
    low_percent = (low_count / total_count) * 100

    # 4. Initialize pure code Presentation with widescreen (16:9) layout
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Helper function for footers
    def add_footer(slide, r_type):
        footer_text = f'CONFIDENTIAL "{r_type or "Internal Audit Executive Summary"}"'
        left_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(6.0), Inches(0.35))
        tf_left = left_box.text_frame
        tf_left.word_wrap = True
        tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0
        p_left = tf_left.paragraphs[0]
        p_left.text = footer_text
        if p_left.runs:
            run_l = p_left.runs[0]
        else:
            run_l = p_left.add_run()
        run_l.font.name = "Inter"
        run_l.font.size = Pt(9)
        run_l.font.italic = True
        run_l.font.color.rgb = RGBColor(148, 163, 184) # light slate-400

        right_box = slide.shapes.add_textbox(Inches(7.0), Inches(6.9), Inches(5.733), Inches(0.4))
        tf_right = right_box.text_frame
        tf_right.word_wrap = True
        tf_right.margin_left = tf_right.margin_right = tf_right.margin_top = tf_right.margin_bottom = 0
        p_right = tf_right.paragraphs[0]
        p_right.alignment = PP_ALIGN.RIGHT
        p_right.text = "Drafted by ALTeX HUB"
        if p_right.runs:
            run_r = p_right.runs[0]
        else:
            run_r = p_right.add_run()
        run_r.font.name = "Inter"
        run_r.font.size = Pt(11)
        run_r.font.bold = True
        run_r.font.color.rgb = RGBColor(71, 85, 105) # slate-600

    # Helper function to add background, top accent bar, slide title, logo, and footer
    def add_slide_decorations(slide, title_text):
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 250, 252)

        # Top Accent bar
        from pptx.enum.shapes import MSO_SHAPE
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(59, 130, 246)
        accent_bar.line.fill.background()

        # Slide Title
        if title_text:
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9.5), Inches(0.8))
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = title_text
            run = p.runs[0]
            run.font.name = "Inter"
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = RGBColor(15, 23, 42)

        # Logo on Top Right
        logo_box = slide.shapes.add_textbox(Inches(10.5), Inches(0.4), Inches(2.233), Inches(0.5))
        tf_logo = logo_box.text_frame
        tf_logo.margin_left = tf_logo.margin_right = tf_logo.margin_top = tf_logo.margin_bottom = 0
        p_logo = tf_logo.paragraphs[0]
        p_logo.alignment = PP_ALIGN.RIGHT
        p_logo.text = "ALTeX HUB"
        run_logo = p_logo.runs[0]
        run_logo.font.name = "Inter"
        run_logo.font.size = Pt(12)
        run_logo.font.bold = True
        run_logo.font.color.rgb = RGBColor(30, 41, 59)

        add_footer(slide, report_type)

    # Slide 1: Cover
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.background
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = RGBColor(15, 23, 42) # Slate 900
    
    # Accent decoration on cover
    from pptx.enum.shapes import MSO_SHAPE
    cover_accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(0.1), Inches(4.5))
    cover_accent.fill.solid()
    cover_accent.fill.fore_color.rgb = RGBColor(59, 130, 246)
    cover_accent.line.fill.background()

    # Cover Text Box
    cover_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(3.8))
    tf_c = cover_box.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_right = tf_c.margin_top = tf_c.margin_bottom = 0
    
    p_sub = tf_c.paragraphs[0]
    p_sub.text = "Internal Audit & Executive Summary"
    r_sub = p_sub.runs[0]
    r_sub.font.name = "Inter"
    r_sub.font.size = Pt(18)
    r_sub.font.bold = True
    r_sub.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
    
    p_title = tf_c.add_paragraph()
    p_title.space_before = Pt(14)
    p_title.text = f"Comprehensive Review: {report_name}"
    r_title = p_title.runs[0]
    r_title.font.name = "Inter"
    r_title.font.size = Pt(36)
    r_title.font.bold = True
    r_title.font.color.rgb = RGBColor(255, 255, 255)

    p_time = tf_c.add_paragraph()
    p_time.space_before = Pt(28)
    p_time.text = f"Timeline: {db_start_date} – {db_end_date}"
    r_time = p_time.runs[0]
    r_time.font.name = "Inter"
    r_time.font.size = Pt(16)
    r_time.font.color.rgb = RGBColor(203, 213, 225) # Slate 300
    
    add_footer(slide1, report_type)

    # Slide 2: Introduction
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide2, "1. Introduction & Background")
    intro_box = slide2.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.133), Inches(5.0))
    tf_intro = intro_box.text_frame
    tf_intro.word_wrap = True
    tf_intro.margin_left = tf_intro.margin_right = tf_intro.margin_top = tf_intro.margin_bottom = 0
    
    p_intro_sub = tf_intro.paragraphs[0]
    p_intro_sub.text = "Context of the Review"
    p_intro_sub.runs[0].font.name = "Inter"
    p_intro_sub.runs[0].font.size = Pt(16)
    p_intro_sub.runs[0].font.bold = True
    p_intro_sub.runs[0].font.color.rgb = RGBColor(30, 41, 59)
    
    bullets = [
        f"In accordance with the 2026 Annual Audit Plan, Internal Audit has initiated a comprehensive review of {company}, a key player in the {sector.lower()} sector.",
        f"Our audit focuses on evaluating the robustness of {company}'s core manufacturing operations, supply chain management, and corresponding financial controls.",
        "The primary objective is to provide independent assurance to the Board regarding the effectiveness of the risk management framework, specifically targeting production efficiency and operational resilience."
    ]
    for b in bullets:
        p_b = tf_intro.add_paragraph()
        p_b.space_before = Pt(14)
        p_b.level = 0
        p_b.text = "•  " + b
        run_b = p_b.runs[0]
        run_b.font.name = "Inter"
        run_b.font.size = Pt(13)
        run_b.font.color.rgb = RGBColor(71, 85, 105)

    # Slide 3: Scope
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide3, "2. Audit Scope")
    scope_box = slide3.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.133), Inches(5.0))
    tf_scope = scope_box.text_frame
    tf_scope.word_wrap = True
    tf_scope.margin_left = tf_scope.margin_right = tf_scope.margin_top = tf_scope.margin_bottom = 0
    
    p_scope_sub = tf_scope.paragraphs[0]
    p_scope_sub.text = "Review Period"
    p_scope_sub.runs[0].font.name = "Inter"
    p_scope_sub.runs[0].font.size = Pt(16)
    p_scope_sub.runs[0].font.bold = True
    p_scope_sub.runs[0].font.color.rgb = RGBColor(30, 41, 59)
    
    p_scope_text = tf_scope.add_paragraph()
    p_scope_text.space_before = Pt(14)
    p_scope_text.text = f"Testing procedures were conducted on transactions, access logs, and operational activities occurring between {db_start_date} and {db_end_date}."
    run_s = p_scope_text.runs[0]
    run_s.font.name = "Inter"
    run_s.font.size = Pt(13)
    run_s.font.color.rgb = RGBColor(71, 85, 105)

    # Slide 4: Methodology
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide4, "3. Approach & Methodology")
    method_box = slide4.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.133), Inches(5.0))
    tf_method = method_box.text_frame
    tf_method.word_wrap = True
    tf_method.margin_left = tf_method.margin_right = tf_method.margin_top = tf_method.margin_bottom = 0
    
    methodologies = [
        ("Risk Assessment", "Conducted preliminary risk scoring based on financial materiality and operational criticality."),
        ("Process Walkthroughs", "Interviewed process owners to document current-state workflows and identify key control points."),
        ("Sample Testing", f"Applied statistical sampling to test the operating effectiveness of {total_count} distinct controls over a 9-month period."),
        ("Evaluation", "Assessed identified control gaps against the COSO framework and evaluated existing compensating controls.")
    ]
    
    first = True
    for title, desc in methodologies:
        p_m = tf_method.paragraphs[0] if first else tf_method.add_paragraph()
        first = False
        if p_m.text:
            p_m.space_before = Pt(14)
        p_m.text = f"•  {title}: "
        r_t = p_m.runs[0]
        r_t.font.name = "Inter"
        r_t.font.size = Pt(13)
        r_t.font.bold = True
        r_t.font.color.rgb = RGBColor(30, 41, 59)
        
        r_d = p_m.add_run()
        r_d.text = desc
        r_d.font.name = "Inter"
        r_d.font.size = Pt(13)
        r_d.font.color.rgb = RGBColor(71, 85, 105)

    # Slide 5: Executive Summary
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide5, "4. Executive Summary")
    exec_box = slide5.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.133), Inches(5.0))
    tf_exec = exec_box.text_frame
    tf_exec.word_wrap = True
    tf_exec.margin_left = tf_exec.margin_right = tf_exec.margin_top = tf_exec.margin_bottom = 0
    
    overall_rating = "Needs Improvement" if high_count > 0 else "Satisfactory"
    
    p_rating = tf_exec.paragraphs[0]
    p_rating.text = f"Based on the results of our testing, the overall control environment is rated as \"{overall_rating}\"."
    r_rat = p_rating.runs[0]
    r_rat.font.name = "Inter"
    r_rat.font.size = Pt(14)
    r_rat.font.bold = True
    r_rat.font.color.rgb = RGBColor(153, 27, 27) if high_count > 0 else RGBColor(22, 101, 52)
    
    p_desc = tf_exec.add_paragraph()
    p_desc.space_before = Pt(14)
    p_desc.text = f"Our audit reviewed {len(controls)} controls across {len(process_counts)} key business processes. We identified {high_count} High Risk findings and {medium_count} Medium Risk findings that require attention."
    r_desc = p_desc.runs[0]
    r_desc.font.name = "Inter"
    r_desc.font.size = Pt(13)
    r_desc.font.color.rgb = RGBColor(71, 85, 105)
    
    p_remed = tf_exec.add_paragraph()
    p_remed.space_before = Pt(14)
    p_remed.text = "Management has proactively acknowledged these findings and has committed resources to execute a robust remediation plan by the end of Q4."
    r_rem = p_remed.runs[0]
    r_rem.font.name = "Inter"
    r_rem.font.size = Pt(13)
    r_rem.font.color.rgb = RGBColor(71, 85, 105)

    # Slide 6: Summary of Key Findings (grouped by process)
    # Group by Process for Slide 6
    process_summary = {}
    for ctrl in controls:
        proc = ctrl['process']
        if proc not in process_summary:
            process_summary[proc] = {
                'process': proc,
                'exception_count': process_counts.get(proc, 0),
                'risk_scores': [],
                'risk_levels': [],
                'auditors': set()
            }
        
        # Risk level sorting: HIGH > MEDIUM > LOW
        if ctrl.get('risk_score'):
            process_summary[proc]['risk_scores'].append(ctrl['risk_score'])
        process_summary[proc]['risk_levels'].append((ctrl.get('risk_level') or 'LOW').upper())
        if ctrl.get('auditor'):
            for aud in ctrl['auditor'].split(','):
                process_summary[proc]['auditors'].add(aud.strip())
                
    final_processes = []
    risk_rank = {"HIGH": 3, "MEDIUM": 2, "LOW": 1}
    for proc_data in process_summary.values():
        scores = proc_data['risk_scores']
        proc_data['risk_score'] = int(sum(scores) / len(scores)) if scores else 0
        proc_data['risk_level'] = max(proc_data['risk_levels'], key=lambda x: risk_rank.get(x, 0)) if proc_data['risk_levels'] else "LOW"
        proc_data['auditor'] = ", ".join(proc_data['auditors'])
        final_processes.append(proc_data)
        
    sorted_processes = sorted(final_processes, key=lambda x: x['exception_count'], reverse=True)

    chunks = [sorted_processes[i:i + 4] for i in range(0, len(sorted_processes), 4)]
    if not chunks:
        chunks = [[]]
        
    for chunk_idx, chunk in enumerate(chunks):
        target_slide = prs.slides.add_slide(blank_layout)
        title_text = "Summary of Key Findings"
        if len(chunks) > 1:
            title_text += f" ({chunk_idx + 1}/{len(chunks)})"
        add_slide_decorations(target_slide, title_text)
            
        # Calculate table height dynamically based on rows
        table_height = 385000 + len(chunk) * 560000
        table_shape = target_slide.shapes.add_table(len(chunk) + 1, 5, 609600, 1420000, 10972800, table_height)
        table = table_shape.table
        
        table.columns[0].width = 2400000
        table.columns[1].width = 2000000
        table.columns[2].width = 1600000
        table.columns[3].width = 1600000
        table.columns[4].width = 3372800
        
        # Apply padding and vertical alignment globally
        for r_idx in range(len(chunk) + 1):
            for c_idx in range(5):
                cell = table.cell(r_idx, c_idx)
                cell.margin_left = Pt(8)
                cell.margin_right = Pt(8)
                cell.margin_top = Pt(4)
                cell.margin_bottom = Pt(4)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        headers = ["Process", "Total Exception Count", "Avg Risk Score", "Risk Level", "Auditor(s)"]
        for col_idx, text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = text
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = "Inter"
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(30, 41, 59)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
        for row_idx, ctrl in enumerate(chunk):
            row = row_idx + 1
            
            # Process
            cell = table.cell(row, 0)
            cell.text = str(ctrl['process'])
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = "Inter"
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = RGBColor(15, 23, 42)
            
            # Exception Count
            cell = table.cell(row, 1)
            cell.text = str(ctrl['exception_count'])
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = "Inter"
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(51, 65, 85)
            
            # Risk Score
            cell = table.cell(row, 2)
            cell.text = str(ctrl['risk_score'])
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = "Inter"
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(51, 65, 85)
            
            # Risk Level
            risk_lvl = ctrl['risk_level'].upper()
            if risk_lvl == "HIGH": risk_color = (153, 27, 27)
            elif risk_lvl == "MEDIUM": risk_color = (146, 64, 14)
            else: risk_color = (22, 101, 52)
            
            cell = table.cell(row, 3)
            cell.text = risk_lvl
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = "Inter"
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = RGBColor(*risk_color)
            
            # Auditor
            cell = table.cell(row, 4)
            cell.text = ctrl['auditor'] or "N/A"
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.name = "Inter"
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(71, 85, 105)

    # Slide 7: Risk Distribution Profile
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide7, "Risk Distribution Profile")
    
    # Load mapping and definitions from ref_audit_plan table in PostgreSQL
    conn_sent = sentinel_pool.getconn()
    try:
        with conn_sent.cursor(cursor_factory=RealDictCursor) as cur:
            cur.execute("SET search_path TO sentinel_db;")
            cur.execute('SELECT * FROM ref_audit_plan;')
            rows_ref = cur.fetchall()
            excel_df = pd.DataFrame(rows_ref)
    finally:
        sentinel_pool.putconn(conn_sent)
    
    table_to_id = {
        "procurement_higher_contract": 1,
        "sjin13_unplanned_delivery": 2,
        "sjpa7_msme_penalty": 3,
        "po_terms_changed": 4,
        "po_split": 5,
        "bank_account_changed": 6,
        "cjsa23_sales_return_qty": 7,
        "sales_return_price_mismatch": 8,
        "multiple_sales_return": 9,
        "sales_return_im": 10,
        "sales_return_180": 11,
        "duplicate_customers": 12,
        "cjsa22_foc_discount": 13,
        "scrap_sales": 14,
        "mjot06_yield_loss": 15,
        "reorder_level": 16,
        "finished_goods_dispatched_wo_qi": 17,
        "cjs1_quality_rejected": 18,
        "tds_insight": 19,
        "gst_working": 20,
        "password_test": 21,
        "direct_changes_sap": 22
    }
    
    id_to_meta = {}
    for idx, row in excel_df.iterrows():
        if pd.isna(row['id']):
            continue
        id_val = int(row['id'])
        id_to_meta[id_val] = {
            "process": row['title'],
            "risk_level": str(row['risk_level']).strip().upper() if not pd.isna(row['risk_level']) else 'LOW',
            "insight_name": row['Insight Name'],
            "risk_description": row['Risk description']
        }
        
    # Re-connect to CompliBear to get the counts for risk profile tables
    conn_comp = psycopg2.connect(**db_config_complibear)
    try:
        with conn_comp.cursor() as cur:
            cur.execute("SET search_path TO complibear;")
            # Fetch existing tables to prevent UndefinedTable exceptions
            cur.execute("SELECT table_name FROM information_schema.tables WHERE table_schema = 'complibear';")
            existing_tables = {r[0] for r in cur.fetchall()}

            table_counts = {}
            for t_name in table_to_id.keys():
                if t_name in existing_tables:
                    cur.execute(f'SELECT COUNT(*) FROM "{t_name}";')
                    table_counts[t_name] = cur.fetchone()[0]
                else:
                    table_counts[t_name] = 0
    finally:
        conn_comp.close()
        
    process_risk_counts = {}
    for t_name, count in table_counts.items():
        id_val = table_to_id[t_name]
        meta = id_to_meta.get(id_val)
        if not meta:
            continue
        proc = meta['process']
        r_level = meta['risk_level']
        
        if proc not in process_risk_counts:
            process_risk_counts[proc] = {"HIGH": 0, "MEDIUM": 0, "LOW": 0}
        process_risk_counts[proc][r_level] += count
        
    sorted_proc_counts = sorted(
        process_risk_counts.items(),
        key=lambda x: sum(x[1].values()),
        reverse=True
    )
    
    # Section A Table: Risk Rating Definitions (Left Half)
    table_shape_a = slide7.shapes.add_table(4, 2, Pt(40), Pt(115), Pt(410), Pt(320))
    table_a = table_shape_a.table
    table_a.columns[0].width = Pt(110)
    table_a.columns[1].width = Pt(300)
    
    headers_a = ["Risk Rating", "Description"]
    for col_idx, text in enumerate(headers_a):
        cell = table_a.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Inter"
        run.font.size = Pt(10.5)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(10)
        cell.margin_right = Pt(10)
        cell.margin_top = Pt(6)
        cell.margin_bottom = Pt(6)
        
    definitions = [
        ("High Risk", "Risks that could seriously compromise the internal control framework, data integrity, and/or operational efficiency. These risks must be addressed with utmost priority."),
        ("Medium Risk", "Risks that could compromise internal controls, data integrity, and/or operational efficiency. These should be addressed after High Risks. Compensating controls may exist but remediation is still required."),
        ("Low Risk", "Risks that do not seriously affect the system in the short term. However, failure to address them may lead to long-term inefficiencies, non-compliance, and weakening of the control framework.")
    ]
    
    colors_a = {
        "High Risk": (239, 68, 68),
        "Medium Risk": (245, 158, 11),
        "Low Risk": (16, 185, 129)
    }
    
    for row_idx, (rating, desc) in enumerate(definitions):
        row = row_idx + 1
        bg_color = RGBColor(248, 250, 252)
        
        cell_rating = table_a.cell(row, 0)
        cell_rating.text = rating
        cell_rating.fill.solid()
        cell_rating.fill.fore_color.rgb = bg_color
        p = cell_rating.text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.name = "Inter"
        run.font.size = Pt(9.5)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*colors_a[rating])
        cell_rating.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell_rating.margin_left = Pt(10)
        cell_rating.margin_right = Pt(10)
        cell_rating.margin_top = Pt(8)
        cell_rating.margin_bottom = Pt(8)
        
        cell_desc = table_a.cell(row, 1)
        cell_desc.text = desc
        cell_desc.fill.solid()
        cell_desc.fill.fore_color.rgb = bg_color
        p = cell_desc.text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.name = "Inter"
        run.font.size = Pt(9.0)
        run.font.color.rgb = RGBColor(71, 85, 105)
        cell_desc.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell_desc.margin_left = Pt(10)
        cell_desc.margin_right = Pt(10)
        cell_desc.margin_top = Pt(8)
        cell_desc.margin_bottom = Pt(8)
        
    # Section B Table: Risk Distribution by Process (Right Half, Top)
    num_processes = len(sorted_proc_counts)
    rows_b = num_processes + 1
    cols_b = 5
    table_shape_b = slide7.shapes.add_table(rows_b, cols_b, Pt(496), Pt(115), Pt(425), Pt(180))
    table_b = table_shape_b.table
    table_b.columns[0].width = Pt(165)
    table_b.columns[1].width = Pt(65)
    table_b.columns[2].width = Pt(65)
    table_b.columns[3].width = Pt(65)
    table_b.columns[4].width = Pt(65)
    
    headers_b = ["Process", "High", "Medium", "Low", "Total"]
    for col_idx, text in enumerate(headers_b):
        cell = table_b.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Inter"
        run.font.size = Pt(9.5)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(6)
        cell.margin_right = Pt(6)
        cell.margin_top = Pt(4)
        cell.margin_bottom = Pt(4)
        
    for row_idx, (proc, levels) in enumerate(sorted_proc_counts):
        row = row_idx + 1
        bg_color = RGBColor(255, 255, 255) if row % 2 == 1 else RGBColor(248, 250, 252)
        
        cell_proc = table_b.cell(row, 0)
        cell_proc.text = proc
        cell_proc.fill.solid()
        cell_proc.fill.fore_color.rgb = bg_color
        p = cell_proc.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Inter"
        run.font.size = Pt(8.5)
        run.font.bold = True
        run.font.color.rgb = RGBColor(30, 41, 59)
        cell_proc.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell_proc.margin_left = Pt(6)
        cell_proc.margin_right = Pt(6)
        
        for level_idx, level_key in enumerate(["HIGH", "MEDIUM", "LOW"]):
            cell_lvl = table_b.cell(row, level_idx + 1)
            cell_lvl.text = str(levels[level_key])
            cell_lvl.fill.solid()
            cell_lvl.fill.fore_color.rgb = bg_color
            p = cell_lvl.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.name = "Inter"
            run.font.size = Pt(8.5)
            if levels[level_key] > 0:
                run.font.bold = True
                if level_key == "HIGH":
                    run.font.color.rgb = RGBColor(239, 68, 68)
                elif level_key == "MEDIUM":
                    run.font.color.rgb = RGBColor(245, 158, 11)
                else:
                    run.font.color.rgb = RGBColor(16, 185, 129)
            else:
                run.font.color.rgb = RGBColor(148, 163, 184)
            cell_lvl.vertical_anchor = MSO_ANCHOR.MIDDLE
            
        total_val = sum(levels.values())
        cell_tot = table_b.cell(row, 4)
        cell_tot.text = str(total_val)
        cell_tot.fill.solid()
        cell_tot.fill.fore_color.rgb = bg_color
        p = cell_tot.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = "Inter"
        run.font.size = Pt(8.5)
        run.font.bold = True
        run.font.color.rgb = RGBColor(30, 41, 59)
        cell_tot.vertical_anchor = MSO_ANCHOR.MIDDLE
        
    # Section C: Clustered Column Chart (Right Half, Bottom)
    abgrad_map = {
        "Procure to Pay": "P2P",
        "Order to Cash": "O2C",
        "Inventory Management": "Inventory",
        "Quality Management": "Quality",
        "Taxation": "Tax",
        "Cyber Security": "Cyber",
        "Scrap Management": "Scrap"
    }
    
    categories = [abgrad_map.get(proc, proc[:10]) for proc, _ in sorted_proc_counts]
    high_vals = [levels["HIGH"] for _, levels in sorted_proc_counts]
    med_vals = [levels["MEDIUM"] for _, levels in sorted_proc_counts]
    low_vals = [levels["LOW"] for _, levels in sorted_proc_counts]
    
    x = np.arange(len(categories))
    width_bar = 0.25
    
    fig, ax = plt.subplots(figsize=(5.9, 2.35))
    ax.bar(x - width_bar, high_vals, width_bar, label='High', color='#ef4444')
    ax.bar(x, med_vals, width_bar, label='Medium', color='#f59e0b')
    ax.bar(x + width_bar, low_vals, width_bar, label='Low', color='#10b981')
    
    ax.set_ylabel('Findings Count', fontsize=8, color='#475569', fontweight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels(categories, fontsize=8, color='#475569', fontweight='bold')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#cbd5e1')
    ax.spines['bottom'].set_color('#cbd5e1')
    ax.tick_params(axis='y', colors='#475569', labelsize=8)
    ax.yaxis.grid(True, linestyle='--', alpha=0.5, color='#cbd5e1')
    ax.set_axisbelow(True)
    
    ax.legend(frameon=False, loc='upper right', fontsize=8, labelcolor='#475569')
    
    fig.patch.set_alpha(0.0)
    ax.patch.set_alpha(0.0)
    plt.tight_layout()
    
    chart_buf = io.BytesIO()
    fig.savefig(chart_buf, format='png', bbox_inches='tight', transparent=True, dpi=300)
    chart_buf.seek(0)
    plt.close(fig)
    
    slide7.shapes.add_picture(chart_buf, Pt(496), Pt(315), Pt(425), Pt(190))

    # Slide 8: Insights/Exceptions by Process
    slide8 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide8, "Insights/Exceptions by Process")

    sorted_process_counts = sorted(process_counts.items(), key=lambda x: x[1])
    sorted_process_counts = [x for x in sorted_process_counts if x[1] > 0]
    proc_names = [x[0] for x in sorted_process_counts]
    proc_vals = [x[1] for x in sorted_process_counts]
    
    fig, ax = plt.subplots(figsize=(11.66, 4.0))
    plt.subplots_adjust(left=0.25, right=0.9, top=0.9, bottom=0.15)
    
    num_cats = len(proc_names)
    bar_height = max(0.2, min(0.6, 10 / max(num_cats, 1)))
    
    bars = ax.barh(proc_names, proc_vals, color='#3b82f6', edgecolor='white', height=bar_height)
    for bar in bars:
        width_val = bar.get_width()
        ax.text(width_val + (width_val * 0.02 if width_val > 0 else 10), 
                bar.get_y() + bar.get_height()/2, 
                f"{width_val:,}", 
                va='center', ha='left', fontsize=9, fontweight='bold', color='#1e293b')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_visible(False)
    ax.spines['left'].set_color('#94a3b8')
    ax.xaxis.set_visible(False)
    
    label_size = 10 if num_cats < 10 else 8
    ax.tick_params(axis='y', colors='#475569', labelsize=label_size)
    
    fig.patch.set_alpha(0.0)
    ax.patch.set_alpha(0.0)
    
    chart_buf2 = io.BytesIO()
    plt.savefig(chart_buf2, format='png', bbox_inches='tight', dpi=180, transparent=True)
    chart_buf2.seek(0)
    plt.close(fig)
    
    slide8.shapes.add_picture(chart_buf2, Inches(0.6), Inches(1.5), Inches(12.133), Inches(4.5))
    
    max_process = max(process_counts, key=process_counts.get) if process_counts else "Unknown"
    desc_box = slide8.shapes.add_textbox(Inches(0.6), Inches(6.1), Inches(12.133), Inches(0.8))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    tf_desc.margin_left = tf_desc.margin_right = tf_desc.margin_top = tf_desc.margin_bottom = 0
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = f"{max_process} represents the largest volume of findings, driving the primary narrative and remediation focus for the remainder of the fiscal year."
    p_desc.runs[0].font.name = "Inter"
    p_desc.runs[0].font.size = Pt(11)
    p_desc.runs[0].font.italic = True
    p_desc.runs[0].font.color.rgb = RGBColor(71, 85, 105)

    # Slide 9: Conclusion & Next Steps
    slide9 = prs.slides.add_slide(blank_layout)
    add_slide_decorations(slide9, "6. Conclusion & Next Steps")
    
    conclusion_box = slide9.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.133), Inches(5.0))
    tf_conclusion = conclusion_box.text_frame
    tf_conclusion.word_wrap = True
    tf_conclusion.margin_left = tf_conclusion.margin_right = tf_conclusion.margin_top = tf_conclusion.margin_bottom = 0
    
    p_conc_sub = tf_conclusion.paragraphs[0]
    p_conc_sub.text = "Path Forward to Compliance"
    p_conc_sub.runs[0].font.name = "Inter"
    p_conc_sub.runs[0].font.size = Pt(16)
    p_conc_sub.runs[0].font.bold = True
    p_conc_sub.runs[0].font.color.rgb = RGBColor(30, 41, 59)
    
    conc_bullets = [
        "The proactive stance taken by executive management to address the IT and vendor management deficiencies highlights a strong organizational commitment to a mature control environment.",
        "Execution: Complete the Q4 remediation plan focusing on automated access de-provisioning protocols.",
        "Validation: Internal Audit to conduct formal validation testing of all remediated controls in January 2027.",
        "Monitoring: Implement continuous monitoring dashboards to track privileged system access requests."
    ]
    for cb in conc_bullets:
        p_cb = tf_conclusion.add_paragraph()
        p_cb.space_before = Pt(14)
        p_cb.level = 0
        p_cb.text = "•  " + cb
        p_cb.runs[0].font.name = "Inter"
        p_cb.runs[0].font.size = Pt(13)
        p_cb.runs[0].font.color.rgb = RGBColor(71, 85, 105)

    # Create Key Observations & Detailed Observation slides if Detailed Findings Report
    if report_type == "Detailed Findings Report":
        # Recommendations mapping
        recommendations = {
            "bank_account_changed": "Implement dual-authorization controls for vendor bank master data changes and callback verification.",
            "cjs1_quality_rejected": "Strengthen SAP gatekeeping controls to block issuing quality-rejected raw materials to production.",
            "cjsa22_foc_discount": "Standardize the approval matrix for Free-of-Cost (FOC) sales and configure automated blocks for unauthorized discounts.",
            "cjsa23_sales_return_qty": "Establish automated batch validation checks to prevent sales return quantities from exceeding original quantity.",
            "direct_changes_sap": "Restrict direct tables edit permissions in SAP production and implement continuous activity logging.",
            "duplicate_customers": "Establish a customer master cleansing process and activate duplicate checks during customer creation.",
            "finished_goods_dispatched_wo_qi": "Enforce validation rules in ERP to prevent dispatching finished goods without approved QI certificate.",
            "gst_working": "Perform monthly automated reconciliations of ITC between GSTR-2B and purchase registers.",
            "mjot06_yield_loss": "Standardize bill of materials (BOM) formulas and perform daily variance analysis on production yield losses.",
            "multiple_sales_return": "Configure system controls to restrict duplicate sales return references against a single invoice.",
            "password_test": "Enforce strong password complexity rules, MFA, and mandatory 90-day password rotation policies.",
            "po_split": "Implement purchase order aggregation controls and monitor split PO behaviors using analytics.",
            "po_terms_changed": "Enforce systemic locks on purchase orders once GRN is posted, requiring VP approval for any post-GRN amendments.",
            "procurement_higher_contract": "Activate system checks in SAP to block purchase orders with prices exceeding contractually agreed rates.",
            "reorder_level": "Automate reorder trigger warnings in ERP based on real-time inventory levels to prevent premature procurement.",
            "sales_return_180": "Enforce system validation to block sales return claims initiated after 180 days from original invoice date.",
            "sales_return_im": "Investigate immediate sales returns (same day) and require formal reasons for all return authorizations.",
            "sales_return_price_mismatch": "Configure SAP price tolerance limits to reject sales return credits priced higher than selling price.",
            "scrap_sales": "Establish competitive bidding and periodic vendor audits for scrap sales, reconciling weighbridge slips.",
            "sjin13_unplanned_delivery": "Enforce strict tolerance thresholds in SAP (max 5%) for unplanned delivery costs, requiring VP approval.",
            "sjpa7_msme_penalty": "Optimize MSME payment workflows to ensure invoice processing and payments occur within statutory 45 days.",
            "tds_insight": "Automate TDS tax rate validation against vendor PAN database to ensure compliance and avoid incorrect withholding."
        }
        
        # Get active observations (count > 0)
        active_obs = []
        for t_name, count in table_counts.items():
            if count > 0:
                id_val = table_to_id[t_name]
                meta = id_to_meta.get(id_val)
                if meta:
                    active_obs.append({
                        "insight_name": meta["insight_name"],
                        "recommendation": recommendations.get(t_name, "Review process validation rules.")
                    })
                    
        # Key Observations Slide splitting (max 4 per slide)
        key_obs_chunks = [active_obs[i:i + 4] for i in range(0, len(active_obs), 4)]
        for chunk_idx, chunk in enumerate(key_obs_chunks):
            target_slide = prs.slides.add_slide(blank_layout)
            title_text = "Key Observations"
            if len(key_obs_chunks) > 1:
                title_text += f" ({chunk_idx + 1}/{len(key_obs_chunks)})"
            add_slide_decorations(target_slide, title_text)
            
            table_height = Pt(30) + len(chunk) * Pt(45)
            table_shape = target_slide.shapes.add_table(len(chunk) + 1, 2, Pt(48), Pt(112), Pt(864), table_height)
            table = table_shape.table
            table.columns[0].width = Pt(240)
            table.columns[1].width = Pt(624)
            
            headers = ["Insight Name", "Recommendation"]
            for col_idx, text in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = text
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.runs[0]
                run.font.name = "Inter"
                run.font.size = Pt(10)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = Pt(10)
                cell.margin_right = Pt(10)
                
            for row_idx, obs in enumerate(chunk):
                row = row_idx + 1
                bg_color = RGBColor(255, 255, 255) if row % 2 == 1 else RGBColor(248, 250, 252)
                
                c_name = table.cell(row, 0)
                c_name.text = str(obs["insight_name"])
                c_name.fill.solid()
                c_name.fill.fore_color.rgb = bg_color
                p = c_name.text_frame.paragraphs[0]
                run = p.runs[0]
                run.font.name = "Inter"
                run.font.size = Pt(9)
                run.font.bold = True
                run.font.color.rgb = RGBColor(30, 41, 59)
                c_name.vertical_anchor = MSO_ANCHOR.MIDDLE
                c_name.margin_left = Pt(10)
                
                c_rec = table.cell(row, 1)
                c_rec.text = str(obs["recommendation"])
                c_rec.fill.solid()
                c_rec.fill.fore_color.rgb = bg_color
                p = c_rec.text_frame.paragraphs[0]
                run = p.runs[0]
                run.font.name = "Inter"
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(71, 85, 105)
                c_rec.vertical_anchor = MSO_ANCHOR.MIDDLE
                c_rec.margin_left = Pt(10)
                
        # Detailed Observations Slide splitting (max 4 per slide)
        detailed_obs = []
        for t_name, count in table_counts.items():
            id_val = table_to_id[t_name]
            meta = id_to_meta.get(id_val)
            if meta:
                detailed_obs.append({
                    "process": meta["process"],
                    "insight_name": meta["insight_name"],
                    "risk_description": meta["risk_description"],
                    "rows": count,
                    "risk_score": meta["risk_level"].capitalize()
                })
                
        process_detailed = {}
        for obs in detailed_obs:
            proc = obs["process"]
            if proc not in process_detailed:
                process_detailed[proc] = []
            process_detailed[proc].append(obs)
            
        for proc, obs_list in process_detailed.items():
            chunks = [obs_list[i:i + 4] for i in range(0, len(obs_list), 4)]
            for chunk_idx, chunk in enumerate(chunks):
                target_slide = prs.slides.add_slide(blank_layout)
                title_text = f"Detailed Observations: {proc}"
                if len(chunks) > 1:
                    title_text += f" ({chunk_idx + 1}/{len(chunks)})"
                add_slide_decorations(target_slide, title_text)
                        
                table_shape = target_slide.shapes.add_table(len(chunk) + 1, 4, Pt(48), Pt(112), Pt(864), Pt(30) + len(chunk) * Pt(45))
                table = table_shape.table
                table.columns[0].width = Pt(200)
                table.columns[1].width = Pt(440)
                table.columns[2].width = Pt(100)
                table.columns[3].width = Pt(124)
                
                headers = ["Insight Name", "Risk Description", "Rows", "Risk Score"]
                for col_idx, text in enumerate(headers):
                    cell = table.cell(0, col_idx)
                    cell.text = text
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER if col_idx >= 2 else PP_ALIGN.LEFT
                    run = p.runs[0]
                    run.font.name = "Inter"
                    run.font.size = Pt(9.5)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    cell.margin_left = Pt(8)
                    cell.margin_right = Pt(8)
                    
                for row_idx, obs in enumerate(chunk):
                    row = row_idx + 1
                    bg_color = RGBColor(255, 255, 255) if row % 2 == 1 else RGBColor(248, 250, 252)
                    
                    cell_ins = table.cell(row, 0)
                    cell_ins.text = str(obs["insight_name"])
                    cell_ins.fill.solid()
                    cell_ins.fill.fore_color.rgb = bg_color
                    p = cell_ins.text_frame.paragraphs[0]
                    run = p.runs[0]
                    run.font.name = "Inter"
                    run.font.size = Pt(8.5)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(30, 41, 59)
                    cell_ins.vertical_anchor = MSO_ANCHOR.MIDDLE
                    cell_ins.margin_left = Pt(8)
                    
                    cell_desc = table.cell(row, 1)
                    cell_desc.text = str(obs["risk_description"])
                    cell_desc.fill.solid()
                    cell_desc.fill.fore_color.rgb = bg_color
                    p = cell_desc.text_frame.paragraphs[0]
                    run = p.runs[0]
                    run.font.name = "Inter"
                    run.font.size = Pt(8.5)
                    run.font.color.rgb = RGBColor(71, 85, 105)
                    cell_desc.vertical_anchor = MSO_ANCHOR.MIDDLE
                    cell_desc.margin_left = Pt(8)
                    
                    cell_rows = table.cell(row, 2)
                    cell_rows.text = f"{obs['rows']:,}"
                    cell_rows.fill.solid()
                    cell_rows.fill.fore_color.rgb = bg_color
                    p = cell_rows.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.runs[0]
                    run.font.name = "Inter"
                    run.font.size = Pt(8.5)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(51, 65, 85)
                    cell_rows.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    risk_lvl = obs["risk_score"].upper()
                    if risk_lvl == "HIGH":
                        risk_color = (239, 68, 68)
                    elif risk_lvl == "MEDIUM":
                        risk_color = (245, 158, 11)
                    else:
                        risk_color = (16, 185, 129)
                        
                    cell_score = table.cell(row, 3)
                    cell_score.text = obs["risk_score"]
                    cell_score.fill.solid()
                    cell_score.fill.fore_color.rgb = bg_color
                    p = cell_score.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    run = p.runs[0]
                    run.font.name = "Inter"
                    run.font.size = Pt(8.5)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(*risk_color)
                    cell_score.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Save to memory buffer
    out_stream = io.BytesIO()
    prs.save(out_stream)
    out_stream.seek(0)
    return out_stream
